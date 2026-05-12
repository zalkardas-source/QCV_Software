import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── Corporate Color Palette ──────────────────────────────────────────
CORP_BLUE      = RGBColor(0, 86, 156)     # #00569c  – primary brand
CORP_BLUE_DARK = RGBColor(0, 58, 107)     # #003a6b  – dark accent
WHITE          = RGBColor(255, 255, 255)
NEAR_BLACK     = RGBColor(30, 30, 35)
DARK_GREY      = RGBColor(65, 65, 75)
MID_GREY       = RGBColor(120, 120, 130)
SKILL_BAR_BG   = RGBColor(225, 230, 238)

# ── Dimensions (Widescreen 13.333 × 7.5 in) ─────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────
def _add_rect(slide, left, top, width, height, fill_rgb, border=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not border:
        shape.line.fill.background()
    return shape

def _add_text(slide, left, top, width, height, text, font_size=11,
              bold=False, italic=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT,
              font_name="Calibri", word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = str(text) if text else ""
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def _add_skill_bar(slide, left, top, width, label, rating, max_rating=10):
    bar_h = Inches(0.18)
    label_w = Inches(2.2)
    bar_w = width - label_w - Inches(0.15)
    clean_label = str(label)[:35] + "..." if len(str(label)) > 38 else str(label)
    _add_text(slide, left, top - Inches(0.02), label_w, Inches(0.25),
              clean_label, font_size=9, color=DARK_GREY, bold=True)
    bar_left = left + label_w + Inches(0.15)
    _add_rect(slide, bar_left, top, bar_w, bar_h, SKILL_BAR_BG)
    try:
        r = int(rating)
    except:
        r = 0
    if r > 0:
        filled_w = bar_w * (min(r, 10) / max_rating)
        _add_rect(slide, bar_left, top, filled_w, bar_h, CORP_BLUE)

def create_pptx_summary(data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    personal = data.get("personal_information", {})
    full_name = personal.get("full_name", "Candidate Profile")
    header_h = Inches(1.35)
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, header_h, CORP_BLUE_DARK)
    
    clean_name = full_name[:40] + "..." if len(full_name) > 43 else full_name
    _add_text(slide, Inches(0.7), Inches(0.2), Inches(8), Inches(0.6),
              clean_name, font_size=30, bold=True, color=WHITE, font_name="Calibri Light")

    contact_parts = []
    if personal.get("email"):    contact_parts.append(f"✉  {personal['email']}")
    if personal.get("phone"):    contact_parts.append(f"✆  {personal['phone']}")
    if personal.get("location"): contact_parts.append(f"⌂  {personal['location']}")
    if contact_parts:
        _add_text(slide, Inches(0.7), Inches(0.82), Inches(10), Inches(0.35),
                  "     |     ".join(contact_parts), font_size=10, color=RGBColor(200, 215, 235))

    _add_text(slide, Inches(9.5), Inches(0.25), Inches(3.5), Inches(0.35),
              "Quatelio CVision", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # ── Executive Summary (Slide 1) ──────────────────────────────
    summary = data.get("small_summary", "")
    section_y = header_h + Inches(0.25)
    if summary:
        _add_text(slide, Inches(0.7), section_y, Inches(3), Inches(0.3),
                  "EXECUTIVE SUMMARY", font_size=11, bold=True, color=CORP_BLUE)
        _add_rect(slide, Inches(0.7), section_y + Inches(0.28), Inches(1.6), Inches(0.025), CORP_BLUE)
        clean_summary = summary[:600] + ("..." if len(summary) > 600 else "")
        _add_text(slide, Inches(0.7), section_y + Inches(0.4), Inches(11.9), Inches(0.8),
                  clean_summary, font_size=9, color=DARK_GREY)
        skill_start_y = section_y + Inches(1.3)
    else:
        skill_start_y = section_y

    # ── Core Skills (Slide 1, 3-Column Layout) ──────────────────
    skill_groups = data.get("skill_matrix", [])
    _add_text(slide, Inches(0.7), skill_start_y, Inches(3), Inches(0.3),
              "CORE SKILLS & COMPETENCES", font_size=11, bold=True, color=CORP_BLUE)
    _add_rect(slide, Inches(0.7), skill_start_y + Inches(0.28), Inches(2.5), Inches(0.025), CORP_BLUE)

    # Dynamic styling based on total skill count
    total_skills = sum(len(g.get("skills", [])) for g in skill_groups if isinstance(g, dict))
    s_font = 8 if total_skills > 30 else 9
    s_line_h = Inches(0.24) if total_skills > 30 else Inches(0.28)
    
    s_y_start = skill_start_y + Inches(0.45)
    s_col_w = Inches(3.9)
    s_cols_x = [Inches(0.7), Inches(4.9), Inches(9.1)]
    
    curr_y = s_y_start
    curr_col = 0
    remaining_groups = []
    # Footer starts at SLIDE_H - Inches(0.35) = 7.15". 
    # We stop at 6.8" to have a safe margin.
    y_limit = Inches(6.8)

    for idx, group in enumerate(skill_groups):
        if not isinstance(group, dict): continue
        if curr_col >= 3:
            remaining_groups = skill_groups[idx:]
            break
        
        # Check if even the category header fits
        if curr_y + Inches(0.4) > y_limit:
            curr_col += 1
            if curr_col >= 3:
                remaining_groups = skill_groups[idx:]
                break
            curr_y = s_y_start

        cat_name = group.get("category", "General").upper()
        _add_text(slide, s_cols_x[curr_col], curr_y, s_col_w, Inches(0.25),
                  cat_name, font_size=s_font - 1, bold=True, color=MID_GREY)
        curr_y += Inches(0.2)

        skills_list = group.get("skills", [])
        for s_idx, s in enumerate(skills_list):
            # Check if next skill fits
            if curr_y + s_line_h > y_limit:
                curr_col += 1
                if curr_col >= 3:
                    remaining_groups = [{"category": group.get("category"), "skills": skills_list[s_idx:]}] + skill_groups[idx+1:]
                    break
                curr_y = s_y_start
            
            _add_skill_bar(slide, s_cols_x[curr_col], curr_y, s_col_w - Inches(0.2), 
                           s.get("skill",""), s.get("rating",5))
            curr_y += s_line_h
        
        if curr_col >= 3: break
        curr_y += Inches(0.15) # Small gap between groups

    # Footer Slide 1
    _add_rect(slide, Emu(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), CORP_BLUE_DARK)

    # ── Skills Continuation Slide (if needed) ──────────────────
    if remaining_groups:
        skill_slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_rect(skill_slide, Emu(0), Emu(0), SLIDE_W, Inches(0.6), CORP_BLUE_DARK)
        _add_text(skill_slide, Inches(0.7), Inches(0.1), Inches(10), Inches(0.4),
                  f"{full_name} - CORE SKILLS (continued)", font_size=16, bold=True, color=WHITE)
        _add_rect(skill_slide, Emu(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), CORP_BLUE_DARK)
        
        # 3-column layout again for continuation
        cs_y_start = Inches(1.0)
        cs_y = cs_y_start
        cs_col = 0
        cs_lines = 0
        cs_max = 18

        for group in remaining_groups:
            if cs_col >= 3: break
            cat_name = group.get("category", "General").upper()
            _add_text(skill_slide, s_cols_x[cs_col], cs_y, s_col_w, Inches(0.25),
                      cat_name, font_size=8, bold=True, color=CORP_BLUE)
            cs_y += Inches(0.25)
            cs_lines += 1

            for s in group.get("skills", []):
                if cs_lines >= cs_max:
                    cs_col += 1
                    if cs_col >= 3: break
                    cs_y = cs_y_start
                    cs_lines = 0
                _add_skill_bar(skill_slide, s_cols_x[cs_col], cs_y, s_col_w - Inches(0.2), 
                               s.get("skill",""), s.get("rating",5))
                cs_y += Inches(0.28)
                cs_lines += 1
            cs_y += Inches(0.15)
            cs_lines += 0.5

    # ── Professional Experience (Slide 2+) ───────────────────────
    projects = data.get("projects", [])
    if projects:
        def add_exp_slide(prs, name, is_cont=False):
            new_s = prs.slides.add_slide(prs.slide_layouts[6])
            _add_rect(new_s, Emu(0), Emu(0), SLIDE_W, Inches(0.6), CORP_BLUE_DARK)
            suffix = " (continued)" if is_cont else ""
            _add_text(new_s, Inches(0.7), Inches(0.1), Inches(10), Inches(0.4),
                      f"{name} - PROFESSIONAL EXPERIENCE{suffix}", font_size=16, bold=True, color=WHITE)
            _add_rect(new_s, Emu(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), CORP_BLUE_DARK)
            return new_s, Inches(0.8)

        curr_exp_slide, curr_exp_y = add_exp_slide(prs, full_name)
        for proj in projects:
            if curr_exp_y > Inches(6.5):
                curr_exp_slide, curr_exp_y = add_exp_slide(prs, full_name, is_cont=True)

            _add_rect(curr_exp_slide, Inches(0.7), curr_exp_y, Inches(0.04), Inches(0.7), CORP_BLUE)
            _add_text(curr_exp_slide, Inches(0.88), curr_exp_y - Inches(0.05), Inches(9.5), Inches(0.3),
                      proj.get("name","Position"), font_size=12, bold=True)
            if proj.get("duration"):
                _add_text(curr_exp_slide, Inches(10.5), curr_exp_y - Inches(0.05), Inches(2), Inches(0.3),
                          proj.get("duration"), font_size=10, color=MID_GREY, align=PP_ALIGN.RIGHT, italic=True)
            if proj.get("description"):
                _add_text(curr_exp_slide, Inches(0.88), curr_exp_y + Inches(0.25), Inches(11.5), Inches(0.5),
                          proj.get("description","")[:500], font_size=10, color=DARK_GREY)
            curr_exp_y += Inches(0.95)

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream.read()
